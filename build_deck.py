"""Generate the Card-D demo PowerPoint from deck_content.py.

Pure python-pptx — native, editable shapes/charts, no headless renderer.
Run: ``python build_deck.py``  ->  ``outputs/total-recall-card-d.pptx``
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import deck_content as C

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
OUT_PATH = os.path.join("outputs", "total-recall-card-d.pptx")


def rgb(name_or_tuple):
    """Accept a BRAND key or an (r,g,b) tuple, return an RGBColor."""
    val = C.BRAND[name_or_tuple] if isinstance(name_or_tuple, str) else name_or_tuple
    return RGBColor(*val)


# --- low-level builders ----------------------------------------------------
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _set_run(run, text, size, color, bold=False, italic=False, font=None):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font or C.BRAND["body_font"]
    run.font.color.rgb = rgb(color)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def _bg(slide, color="white"):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(rect, color)
    rect.shadow.inherit = False
    # send to back so other shapes render on top
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return rect


def _accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    _fill(bar, "accent")
    bar.shadow.inherit = False


def _heading(slide, title, kicker=None):
    """Standard slide heading; returns the bottom Y to start body content."""
    _accent_bar(slide)
    top = Inches(0.45)
    if kicker:
        _, tf = _textbox(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.4))
        _set_run(tf.paragraphs[0].add_run(), kicker.upper(), 13, "accent",
                 bold=True, font=C.BRAND["title_font"])
        top = Inches(0.85)
    box, tf = _textbox(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9))
    _set_run(tf.paragraphs[0].add_run(), title, 30, "ink", bold=True,
             font=C.BRAND["title_font"])
    return top + Inches(1.0)


def _bullets(slide, items, left, top, width, height, size=16,
             color="ink", bullet_color="accent", gap=8):
    """items: list of strings or (text, bold) tuples."""
    _, tf = _textbox(slide, left, top, width, height)
    for i, item in enumerate(items):
        text, bold = (item, False) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        _set_run(r, "•  ", size, bullet_color, bold=True)
        r2 = p.add_run()
        _set_run(r2, text, size, color, bold=bold)
    return tf


def _panel(slide, left, top, width, height, fill="surface"):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(panel, fill)
    panel.shadow.inherit = False
    return panel


# --- slide templates -------------------------------------------------------
def title_slide(prs):
    s = _blank_slide(prs)
    _bg(s, "ink")
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.85), SLIDE_W, Inches(0.12))
    _fill(band, "accent")
    _, tf = _textbox(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(1.6))
    _set_run(tf.paragraphs[0].add_run(), "TOTAL RECALL", 54, "white", bold=True,
             font=C.BRAND["title_font"])
    p = tf.add_paragraph()
    _set_run(p.add_run(), "An agent that visibly gets sharper across sessions",
             24, "accent", font=C.BRAND["title_font"])
    _, tf2 = _textbox(s, MARGIN, Inches(5.1), SLIDE_W - 2 * MARGIN, Inches(1.4))
    _set_run(tf2.paragraphs[0].add_run(),
             "Card D — Sales Engineer for Helios", 20, "white", bold=True)
    p2 = tf2.add_paragraph()
    _set_run(p2.add_run(),
             "Claude Managed Agents + the Memory tool   •   Prospect: Northwind Bank",
             14, "muted")
    return s


def section_quote_slide(prs, kicker, big_text, sub=None, fill="white",
                        big_color="ink", accent_word_color="accent"):
    s = _blank_slide(prs)
    _bg(s, fill)
    _accent_bar(s)
    _, tf = _textbox(s, MARGIN, Inches(0.6), SLIDE_W - 2 * MARGIN, Inches(0.5))
    _set_run(tf.paragraphs[0].add_run(), kicker.upper(), 14, "accent", bold=True,
             font=C.BRAND["title_font"])
    box, tf2 = _textbox(s, MARGIN, Inches(2.3), SLIDE_W - 2 * MARGIN, Inches(2.6))
    tf2.word_wrap = True
    _set_run(tf2.paragraphs[0].add_run(), big_text, 34, big_color, bold=True,
             font=C.BRAND["title_font"])
    if sub:
        _, tf3 = _textbox(s, MARGIN, Inches(5.0), SLIDE_W - 2 * MARGIN, Inches(1.6))
        _set_run(tf3.paragraphs[0].add_run(), sub, 18, "muted")
    return s


def bullets_slide(prs, title, items, kicker=None, intro=None, size=18):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, title, kicker)
    if intro:
        _, tf = _textbox(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.7))
        _set_run(tf.paragraphs[0].add_run(), intro, 16, "muted", italic=True)
        top += Inches(0.8)
    _bullets(s, items, MARGIN, top, SLIDE_W - 2 * MARGIN,
             SLIDE_H - top - Inches(0.4), size=size)
    return s


def why_memory_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    _heading(s, "Memory is not a vector database", kicker="Why this matters")
    # left: the misconception, right: what it really is
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    lp = _panel(s, MARGIN, Inches(2.0), col_w, Inches(4.6), "surface")
    _, tf = _textbox(s, MARGIN + Inches(0.3), Inches(2.3), col_w - Inches(0.6), Inches(4.0))
    _set_run(tf.paragraphs[0].add_run(), "What clients think it is", 18, "muted", bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(10)
    _set_run(p.add_run(), "“A vector database for documents.”", 20, "ink", italic=True)
    p = tf.add_paragraph(); p.space_before = Pt(10)
    _set_run(p.add_run(), "Store everything, retrieve by similarity, hope the right chunk comes back.",
             15, "muted")

    rx = MARGIN + col_w + Inches(0.4)
    rp = _panel(s, rx, Inches(2.0), col_w, Inches(4.6), "ink")
    _, tf2 = _textbox(s, rx + Inches(0.3), Inches(2.3), col_w - Inches(0.6), Inches(4.0))
    _set_run(tf2.paragraphs[0].add_run(), "What it actually is", 18, "accent", bold=True)
    p = tf2.add_paragraph(); p.space_before = Pt(10)
    _set_run(p.add_run(),
             "The agent decides what to remember, what to forget, and what to update when it learns something new.",
             20, "white", bold=True)
    for line in ["Remember the durable facts.", "Forget the noise.",
                 "Update when new info contradicts the old."]:
        p = tf2.add_paragraph(); p.space_before = Pt(8)
        r = p.add_run(); _set_run(r, "•  ", 16, "accent", bold=True)
        r2 = p.add_run(); _set_run(r2, line, 16, "white")
    return s


def scenario_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "The scenario", kicker="Card D — Helios Sales Engineer")
    _, tf = _textbox(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.0))
    _set_run(tf.paragraphs[0].add_run(), C.SCENARIO["product"] + "   •   ", 16, "ink", bold=True)
    _set_run(tf.paragraphs[0].add_run(), "Prospect: " + C.SCENARIO["prospect"], 16, "muted")
    p = tf.add_paragraph(); p.space_before = Pt(6)
    _set_run(p.add_run(), C.SCENARIO["persona"], 14, "muted", italic=True)

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    rtop = top + Inches(1.5)
    ph = Inches(3.2)
    _panel(s, MARGIN, rtop, col_w, ph, "surface")
    _, t1 = _textbox(s, MARGIN + Inches(0.25), rtop + Inches(0.15), col_w - Inches(0.5), ph - Inches(0.3))
    _set_run(t1.paragraphs[0].add_run(), "ROUND 1  —  what we know", 14, "session1", bold=True)
    for name, desc in C.SCENARIO["round1_docs"]:
        p = t1.add_paragraph(); p.space_before = Pt(8)
        _set_run(p.add_run(), name + " ", 13, "ink", bold=True)
        p2 = t1.add_paragraph()
        _set_run(p2.add_run(), desc, 11.5, "muted")

    rx = MARGIN + col_w + Inches(0.4)
    _panel(s, rx, rtop, col_w, ph, "surface")
    _, t2 = _textbox(s, rx + Inches(0.25), rtop + Inches(0.15), col_w - Inches(0.5), ph - Inches(0.3))
    _set_run(t2.paragraphs[0].add_run(), "ROUND 2  —  what's new (contradicts)", 14, "session2", bold=True)
    for name, desc in C.SCENARIO["round2_docs"]:
        p = t2.add_paragraph(); p.space_before = Pt(8)
        _set_run(p.add_run(), name + " ", 13, "ink", bold=True)
        p2 = t2.add_paragraph()
        _set_run(p2.add_run(), desc, 11.5, "muted")
    return s


def _arrow(slide, x1, y1, x2, y2, color="muted"):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight connector
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(2.25)
    try:
        conn.line.headEnd = "none"
    except Exception:
        pass
    return conn


def architecture_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    _heading(s, "How it works", kicker="Architecture")

    def box(left, top, w, h, label, sub, fill, txt="white"):
        shp = _panel(s, left, top, w, h, fill)
        _, tf = _textbox(s, left + Inches(0.12), top + Inches(0.12), w - Inches(0.24), h - Inches(0.24))
        tf.word_wrap = True
        _set_run(tf.paragraphs[0].add_run(), label, 15, txt, bold=True)
        if sub:
            p = tf.add_paragraph()
            _set_run(p.add_run(), sub, 11, txt)
        return shp

    row_y = Inches(2.4)
    bw, bh = Inches(2.7), Inches(1.5)
    gap = Inches(0.55)
    x0 = MARGIN + Inches(0.1)
    b1 = box(x0, row_y, bw, bh, "Sales Engineer", "asks the question", "session1")
    b2 = box(x0 + (bw + gap), row_y, bw, bh, "Managed Agent", "Claude + Memory tool", "accent")
    b3 = box(x0 + 2 * (bw + gap), row_y, bw, bh, "Memory Store",
             "path-organized\n/customers /objections /competitive /pitch", "ink")
    _arrow(s, x0 + bw, row_y + bh / 2, x0 + bw + gap, row_y + bh / 2, "muted")
    _arrow(s, x0 + 2 * bw + gap, row_y + bh / 2, x0 + 2 * (bw + gap), row_y + bh / 2, "muted")

    # docs feeding in
    docs = box(x0 + (bw + gap), row_y + bh + Inches(0.7), bw, Inches(1.1),
               "Files API → round1 / round2 docs", "uploaded per session", "session2")
    _arrow(s, x0 + (bw + gap) + bw / 2, row_y + bh + Inches(0.7),
           x0 + (bw + gap) + bw / 2, row_y + bh, "muted")

    # footnote: two kinds of memory
    _, tf = _textbox(s, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(1.1))
    _set_run(tf.paragraphs[0].add_run(), "Two kinds of memory:  ", 14, "accent", bold=True)
    _set_run(tf.paragraphs[0].add_run(),
             "in-session (the context window + built-in compaction) vs across-session "
             "(the persistent Memory Store). Sessions come and go; the store endures.",
             14, "muted")
    return s


def session_answer_slide(prs, n, data, color):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "Session %d — the answer" % n,
                   kicker="Q: " + C.TEST_QUESTION)
    # preface panel
    _panel(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.85), "surface")
    _, tf = _textbox(s, MARGIN + Inches(0.2), top + Inches(0.1),
                     SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.65))
    _set_run(tf.paragraphs[0].add_run(), data["preface"], 13, "muted", italic=True)
    # headline
    htop = top + Inches(1.05)
    _, tf2 = _textbox(s, MARGIN, htop, SLIDE_W - 2 * MARGIN, Inches(0.7))
    _set_run(tf2.paragraphs[0].add_run(), data["headline"], 20, color, bold=True)
    # bullets
    _bullets(s, data["bullets"], MARGIN, htop + Inches(0.8),
             SLIDE_W - 2 * MARGIN, Inches(2.9), size=14, bullet_color=color)
    # gap note
    _, tf3 = _textbox(s, MARGIN, Inches(6.7), SLIDE_W - 2 * MARGIN, Inches(0.6))
    _set_run(tf3.paragraphs[0].add_run(), "▸ " + data["gap"], 13, "accent_dark", bold=True)
    return s


def between_calls_slide(prs):
    items = []
    for name, desc in C.SCENARIO["round2_docs"]:
        items.append((name, True))
        items.append(desc)
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "Between the calls", kicker="Two new round-2 inputs")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    ph = Inches(3.6)
    titles = ["NEW OBJECTION (latest call)", "COMPETITIVE UPDATE (from product)"]
    for i, (name, desc) in enumerate(C.SCENARIO["round2_docs"]):
        x = MARGIN + i * (col_w + Inches(0.4))
        _panel(s, x, top, col_w, ph, "surface")
        _, tf = _textbox(s, x + Inches(0.25), top + Inches(0.2), col_w - Inches(0.5), ph - Inches(0.4))
        _set_run(tf.paragraphs[0].add_run(), titles[i], 14, "session2", bold=True)
        p = tf.add_paragraph(); p.space_before = Pt(8)
        _set_run(p.add_run(), name, 14, "ink", bold=True)
        p2 = tf.add_paragraph(); p2.space_before = Pt(8)
        _set_run(p2.add_run(), desc, 15, "ink")
    _, tf = _textbox(s, MARGIN, top + ph + Inches(0.25), SLIDE_W - 2 * MARGIN, Inches(0.7))
    _set_run(tf.paragraphs[0].add_run(),
             "These contradict / extend round 1 — the agent must reconcile, not just append.",
             15, "muted", italic=True)
    return s


def comparison_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    _heading(s, "Same question, sharper answer", kicker="The headline")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    top = Inches(2.0)
    ph = Inches(4.9)
    # session 1 column
    _panel(s, MARGIN, top, col_w, ph, "surface")
    hdr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, top, col_w, Inches(0.6))
    _fill(hdr1, "session1")
    _, h1 = _textbox(s, MARGIN + Inches(0.2), top + Inches(0.06), col_w - Inches(0.4), Inches(0.5))
    _set_run(h1.paragraphs[0].add_run(), "SESSION 1  —  cold start", 15, "white", bold=True)
    _bullets(s, [C.SESSION_1_ANSWER["headline"]] + C.SESSION_1_ANSWER["bullets"][:4],
             MARGIN + Inches(0.2), top + Inches(0.8), col_w - Inches(0.4), ph - Inches(1.0),
             size=12.5, bullet_color="session1")

    rx = MARGIN + col_w + Inches(0.4)
    _panel(s, rx, top, col_w, ph, "surface")
    hdr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, top, col_w, Inches(0.6))
    _fill(hdr2, "session2")
    _, h2 = _textbox(s, rx + Inches(0.2), top + Inches(0.06), col_w - Inches(0.4), Inches(0.5))
    _set_run(h2.paragraphs[0].add_run(), "SESSION 2  —  with memory + new context", 15, "white", bold=True)
    _bullets(s, [C.SESSION_2_ANSWER["headline"]] + C.SESSION_2_ANSWER["bullets"][:4],
             rx + Inches(0.2), top + Inches(0.8), col_w - Inches(0.4), ph - Inches(1.0),
             size=12.5, bullet_color="session2")
    return s


def memory_diff_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "What the agent learned between calls",
                   kicker="S8 — Memory diff")
    rows = []
    for path, content in C.MEMORY_DIFF["added"]:
        rows.append(("ADDED", path, content, "added"))
    for path, content in C.MEMORY_DIFF["changed"]:
        rows.append(("CHANGED", path, content, "changed"))
    n_unchanged = len(C.MEMORY_DIFF["unchanged"])

    n = len(rows) + 1  # +header
    table_w = SLIDE_W - 2 * MARGIN
    th = Inches(4.5)
    tbl_shape = s.shapes.add_table(n + 1, 3, MARGIN, top, table_w, th)
    table = tbl_shape.table
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = table_w - Inches(5.0)

    def cell(r, c, text, color="ink", bold=False, size=12, fill=None):
        cl = table.cell(r, c)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill:
            cl.fill.solid(); cl.fill.fore_color.rgb = rgb(fill)
        else:
            cl.fill.solid(); cl.fill.fore_color.rgb = rgb("white")
        tf = cl.text_frame; tf.word_wrap = True
        tf.paragraphs[0].clear()
        _set_run(tf.paragraphs[0].add_run(), text, size, color, bold=bold)

    cell(0, 0, "CHANGE", "white", True, 12, "ink")
    cell(0, 1, "MEMORY PATH", "white", True, 12, "ink")
    cell(0, 2, "CONTENT", "white", True, 12, "ink")
    for i, (tag, path, content, ckey) in enumerate(rows, start=1):
        cell(i, 0, tag, ckey, True, 11)
        cell(i, 1, path, "ink", False, 11)
        cell(i, 2, content, "muted", False, 11)
    # final summary row
    cell(n, 0, "UNCHANGED", "unchanged", True, 11)
    cell(n, 1, "%d entries" % n_unchanged, "unchanged", False, 11)
    cell(n, 2, "carried over from session 1 untouched", "muted", True, 11)

    _, tf = _textbox(s, MARGIN, top + th + Inches(0.15), SLIDE_W - 2 * MARGIN, Inches(0.5))
    _set_run(tf.paragraphs[0].add_run(),
             "“Here's exactly what the agent learned” — added the failover objection, "
             "competitive intel, and the new CISO; updated residency, cost, and the pitch sequence.",
             13, "muted", italic=True)
    return s


def memory_growth_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "Memory grows where it should", kicker="Facts remembered per category")
    chart_data = CategoryChartData()
    chart_data.categories = [c for c, _, _ in C.MEMORY_GROWTH]
    chart_data.add_series("After session 1", tuple(a for _, a, _ in C.MEMORY_GROWTH))
    chart_data.add_series("After session 2", tuple(b for _, _, b in C.MEMORY_GROWTH))
    gf = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, top,
        SLIDE_W - 2 * MARGIN, Inches(4.4), chart_data)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = rgb("session1")
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = rgb("session2")
    _, tf = _textbox(s, MARGIN, top + Inches(4.5), SLIDE_W - 2 * MARGIN, Inches(0.5))
    _set_run(tf.paragraphs[0].add_run(),
             "New competitive intel and a new objection appear; pitch strategy is one "
             "entry that gets rewritten in place rather than duplicated.",
             13, "muted", italic=True)
    return s


def business_value_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "Why enterprise buyers care", kicker="The questions every client asks")
    y = top
    rh = Inches(0.85)
    for q, a in C.BUSINESS_VALUE:
        _panel(s, MARGIN, y, SLIDE_W - 2 * MARGIN, rh - Inches(0.12), "surface")
        _, tf = _textbox(s, MARGIN + Inches(0.25), y + Inches(0.02),
                         SLIDE_W - 2 * MARGIN - Inches(0.5), rh - Inches(0.2))
        _set_run(tf.paragraphs[0].add_run(), q + "   ", 15, "ink", bold=True)
        _set_run(tf.paragraphs[0].add_run(), a, 14, "muted")
        y += rh
    return s


def stretch_overview_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    _heading(s, "Nine ways to make memory real", kicker="Stretch goals S1–S9")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 2
    positions = [(MARGIN, Inches(1.9)), (MARGIN + col_w + Inches(0.6), Inches(1.9)),
                 (MARGIN, Inches(4.45)), (MARGIN + col_w + Inches(0.6), Inches(4.45))]
    ph = Inches(2.35)
    for (tier_title, goals), (x, y) in zip(C.STRETCH_TIERS, positions):
        _panel(s, x, y, col_w, ph, "surface")
        _, tf = _textbox(s, x + Inches(0.25), y + Inches(0.15), col_w - Inches(0.5), ph - Inches(0.3))
        _set_run(tf.paragraphs[0].add_run(), tier_title, 14, "accent_dark", bold=True)
        for gid, gtitle, _what, _why in goals:
            p = tf.add_paragraph(); p.space_before = Pt(6)
            _set_run(p.add_run(), gid + "  ", 13, "accent", bold=True)
            _set_run(p.add_run(), gtitle, 13, "ink", bold=True)
    return s


def tier_detail_slide(prs, tier_index):
    tier_title, goals = C.STRETCH_TIERS[tier_index]
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, tier_title, kicker="Stretch goals")
    # one panel per goal stacked
    avail = SLIDE_H - top - Inches(0.4)
    gh = (avail - Inches(0.3) * (len(goals) - 1)) / len(goals)
    y = top
    for gid, gtitle, what, why in goals:
        _panel(s, MARGIN, y, SLIDE_W - 2 * MARGIN, gh, "surface")
        _, tf = _textbox(s, MARGIN + Inches(0.3), y + Inches(0.15),
                         SLIDE_W - 2 * MARGIN - Inches(0.6), gh - Inches(0.3))
        _set_run(tf.paragraphs[0].add_run(), gid + " — " + gtitle, 17, "accent_dark", bold=True)
        p = tf.add_paragraph(); p.space_before = Pt(5)
        _set_run(p.add_run(), what, 13.5, "ink")
        p2 = tf.add_paragraph(); p2.space_before = Pt(4)
        _set_run(p2.add_run(), "Why it lands:  ", 13, "accent", bold=True)
        _set_run(p2.add_run(), why, 13, "muted", italic=True)
        y += gh + Inches(0.3)
    return s


def how_to_run_slide(prs):
    s = _blank_slide(prs)
    _bg(s)
    top = _heading(s, "Run it end-to-end", kicker="Verification")
    _panel(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.6), "ink")
    _, tf = _textbox(s, MARGIN + Inches(0.35), top + Inches(0.25),
                     SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(4.1))
    for i, line in enumerate(C.HOW_TO_RUN):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        _set_run(p.add_run(), "$ ", 14, "accent", bold=True, font="Consolas")
        _set_run(p.add_run(), line, 13.5, "white", font="Consolas")
    return s


def closing_slide(prs):
    s = _blank_slide(prs)
    _bg(s, "ink")
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), SLIDE_W, Inches(0.1))
    _fill(band, "accent")
    _, tf = _textbox(s, MARGIN, Inches(2.2), SLIDE_W - 2 * MARGIN, Inches(1.2))
    _set_run(tf.paragraphs[0].add_run(),
             "Same question. Two sessions. A visibly sharper answer.",
             30, "white", bold=True, font=C.BRAND["title_font"])
    _, tf2 = _textbox(s, MARGIN, Inches(3.7), SLIDE_W - 2 * MARGIN, Inches(1.4))
    _set_run(tf2.paragraphs[0].add_run(),
             "That's institutional memory: the agent decides what to keep, what to drop, "
             "and what to rewrite when the world changes.",
             18, "muted")
    p = tf2.add_paragraph(); p.space_before = Pt(14)
    _set_run(p.add_run(), "Regenerate this deck anytime:  python build_deck.py",
             14, "accent", font="Consolas")
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)                                                   # 1
    why_memory_slide(prs)                                              # 2
    section_quote_slide(prs, "The demo in one line",
                        "Same question, asked twice. The second answer is sharper "
                        "because the agent remembered — and reconciled what changed.",
                        sub="No new prompt engineering between calls. Just memory.")  # 3
    scenario_slide(prs)                                                # 4
    architecture_slide(prs)                                            # 5
    section_quote_slide(prs, "The test question (both sessions)",
                        "“" + C.TEST_QUESTION + "”",
                        sub="Asked identically in session 1 and session 2.",
                        big_color="accent")                            # 6
    session_answer_slide(prs, 1, C.SESSION_1_ANSWER, "session1")       # 7
    between_calls_slide(prs)                                           # 8
    session_answer_slide(prs, 2, C.SESSION_2_ANSWER, "session2")       # 9
    comparison_slide(prs)                                              # 10
    memory_diff_slide(prs)                                             # 11
    memory_growth_slide(prs)                                           # 12
    business_value_slide(prs)                                          # 13
    stretch_overview_slide(prs)                                        # 14
    tier_detail_slide(prs, 0)                                          # 15
    tier_detail_slide(prs, 1)                                          # 16
    tier_detail_slide(prs, 2)                                          # 17
    tier_detail_slide(prs, 3)                                          # 18
    how_to_run_slide(prs)                                              # 19
    closing_slide(prs)                                                 # 20

    os.makedirs("outputs", exist_ok=True)
    prs.save(OUT_PATH)
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    count = build()
    print("Wrote %s  (%d slides)" % (OUT_PATH, count))
