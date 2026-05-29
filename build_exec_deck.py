"""Generate the 5-slide *executive summary* of the Card-D demo.

Reuses the slide builders, helpers, and content from build_deck.py — this file
only chooses which slides make the executive cut and adds a tailored title and
closing/ask. Run: ``python build_exec_deck.py``
  ->  ``outputs/total-recall-card-d-exec.pptx``
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

import deck_content as C
import build_deck as bd

OUT_PATH = os.path.join("outputs", "total-recall-card-d-exec.pptx")


def exec_title_slide(prs):
    s = bd._blank_slide(prs)
    bd._bg(s, "ink")
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.7), bd.SLIDE_W, Inches(0.12))
    bd._fill(band, "accent")
    _, tf = bd._textbox(s, bd.MARGIN, Inches(1.7), bd.SLIDE_W - 2 * bd.MARGIN, Inches(1.7))
    bd._set_run(tf.paragraphs[0].add_run(), "TOTAL RECALL", 50, "white", bold=True,
                font=C.BRAND["title_font"])
    p = tf.add_paragraph()
    bd._set_run(p.add_run(), "Executive summary", 24, "accent",
                font=C.BRAND["title_font"])
    _, tf2 = bd._textbox(s, bd.MARGIN, Inches(4.95), bd.SLIDE_W - 2 * bd.MARGIN, Inches(1.6))
    bd._set_run(tf2.paragraphs[0].add_run(),
                "An agent that gets sharper across sessions — it remembers, reconciles, "
                "and answers better.", 20, "white", bold=True)
    p2 = tf2.add_paragraph()
    bd._set_run(p2.add_run(),
                "Claude Managed Agents + Memory   •   Card D: Helios Sales Engineer   •   "
                "Prospect: Northwind Bank", 13, "muted")
    return s


def exec_what_why_slide(prs):
    """Condensed 'what it is + why it's different' on one slide."""
    s = bd._blank_slide(prs)
    bd._bg(s)
    top = bd._heading(s, "Memory is not a vector database", kicker="What we built")
    _, tf = bd._textbox(s, bd.MARGIN, top, bd.SLIDE_W - 2 * bd.MARGIN, Inches(1.0))
    bd._set_run(tf.paragraphs[0].add_run(),
                "It means the agent decides what to remember, what to forget, and what to "
                "update when it learns something new — across sessions, on the same account.",
                18, "ink", bold=True)
    # three-up value tiles
    tiles = [
        ("Remembers", "Durable facts about the prospect, every objection, and our best answer."),
        ("Reconciles", "When new info contradicts the old, it updates in place — not blindly appends."),
        ("Proves it", "An audit trail shows exactly what changed between calls."),
    ]
    col_w = (bd.SLIDE_W - 2 * bd.MARGIN - Inches(0.8)) / 3
    y = top + Inches(1.4)
    ph = Inches(2.9)
    for i, (h, body) in enumerate(tiles):
        x = bd.MARGIN + i * (col_w + Inches(0.4))
        bd._panel(s, x, y, col_w, ph, "surface")
        _, t = bd._textbox(s, x + Inches(0.25), y + Inches(0.25), col_w - Inches(0.5), ph - Inches(0.5))
        bd._set_run(t.paragraphs[0].add_run(), h, 20, "accent_dark", bold=True)
        p = t.add_paragraph(); p.space_before = Pt(10)
        bd._set_run(p.add_run(), body, 14, "muted")
    return s


def exec_ask_slide(prs):
    s = bd._blank_slide(prs)
    bd._bg(s, "ink")
    bd._accent_bar(s)
    _, tf = bd._textbox(s, bd.MARGIN, Inches(0.7), bd.SLIDE_W - 2 * bd.MARGIN, Inches(0.5))
    bd._set_run(tf.paragraphs[0].add_run(), "WHY IT MATTERS", 14, "accent", bold=True,
                font=C.BRAND["title_font"])
    _, tf2 = bd._textbox(s, bd.MARGIN, Inches(1.5), bd.SLIDE_W - 2 * bd.MARGIN, Inches(0.8))
    bd._set_run(tf2.paragraphs[0].add_run(),
                "It answers the questions every enterprise buyer asks", 26, "white", bold=True,
                font=C.BRAND["title_font"])
    y = Inches(2.7)
    rh = Inches(0.78)
    for q, a in C.BUSINESS_VALUE:
        _, tf = bd._textbox(s, bd.MARGIN, y, bd.SLIDE_W - 2 * bd.MARGIN, rh)
        r = tf.paragraphs[0].add_run(); bd._set_run(r, "•  ", 15, "accent", bold=True)
        bd._set_run(tf.paragraphs[0].add_run(), q + "   ", 15, "white", bold=True)
        bd._set_run(tf.paragraphs[0].add_run(), a, 14, "muted")
        y += rh
    return s


def build():
    prs = Presentation()
    prs.slide_width = bd.SLIDE_W
    prs.slide_height = bd.SLIDE_H

    exec_title_slide(prs)                       # 1 — title
    exec_what_why_slide(prs)                    # 2 — what it is / why it's different
    bd.comparison_slide(prs)                    # 3 — same question, sharper answer (the proof)
    bd.memory_diff_slide(prs)                   # 4 — exactly what it learned (audit trail)
    exec_ask_slide(prs)                         # 5 — why it matters / buyer questions

    os.makedirs("outputs", exist_ok=True)
    prs.save(OUT_PATH)
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    count = build()
    print("Wrote %s  (%d slides)" % (OUT_PATH, count))
